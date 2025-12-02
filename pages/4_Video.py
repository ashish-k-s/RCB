# pip install python-pptx moviepy pyttsx3
import shutil
import streamlit as st
from pptx import Presentation
from pdf2image import convert_from_path
import subprocess
import wave
from moviepy import VideoFileClip, AudioFileClip, ImageClip, concatenate_videoclips
import os
import glob
from pathlib import Path

from rcb_audio import generate_audio_file_from_transcript
from rcb_init import init_audio_page, init_audio_prompts, init_audio_vars, init_page

st.set_page_config(
    page_title="Video with RCB"
)

st.title("Create Video using RCB")
if 'current_page' not in st.session_state:
    st.session_state.current_page = "Video"
st.session_state.current_page = "Video"

init_page()
init_audio_page()
init_audio_vars()
init_audio_prompts()


# audio_file_name_str = "rcb_generated_audio"
# if 'audio_file_path_wav' not in st.session_state:
#     st.session_state.audio_file_path_wav = '{st.session_state.video_data_dir}/' + audio_file_name_str + '.wav'
# if 'audio_file_path_mp3' not in st.session_state:
#     st.session_state.audio_file_path_mp3 = '{st.session_state.video_data_dir}/' + audio_file_name_str + '.mp3'
# if 'audio_file_path_txt' not in st.session_state:
#     st.session_state.audio_file_path_txt = '{st.session_state.video_data_dir}/' + audio_file_name_str + '.txt'
if 'video_data_dir' not in st.session_state:
    st.session_state.video_data_dir = f"{st.session_state.user_dir}/video"
if 'video_file_path' not in st.session_state:
    st.session_state.video_file_path = f"{st.session_state.video_data_dir}/rcb_generated_video.mp4"
if 'progress_logs' not in st.session_state:
    st.session_state.progress_logs = st.empty()

def clear_old_files():
    """
    Clear old files in the video data directory
    """
    if os.path.exists(st.session_state.video_data_dir):
        shutil.rmtree(st.session_state.video_data_dir)
    os.makedirs(st.session_state.video_data_dir, exist_ok=True)
    st.session_state.progress_logs.info("Cleared old files in video data directory.")
    
def simple_video_creator():
    """
    Simple function to create video from images and audio files
    """
    st.session_state.progress_logs.info("Creating video from images and audio files... This may take a while, please wait.")

    print(f"\nCreating video from images and audio files in directory: {st.session_state.video_data_dir}")
    # Get image files
    pngs = glob.glob(f"{st.session_state.video_data_dir}/*.png")
    pngs.sort()

    # Get audio files
    mp3s = glob.glob(f"{st.session_state.video_data_dir}/*.mp3")
    mp3s.sort()
     
    # Print file names for debugging
    print("\nImage files:")
    for png in pngs:
        print(f"  {Path(png).name}")
    
    print("\nAudio files:")
    for mp3 in mp3s:
        print(f"  {Path(mp3).name}")
    
    # Create video clips
    clips = []
    
    # Process pairs of images and audio
    pairs = min(len(pngs), len(mp3s))
    for i in range(pairs):
        print(f"\nProcessing pair {i+1}/{pairs}")

        img_path = pngs[i]
        aud_path = mp3s[i]
        
        print(f"  Image: {Path(img_path).name}")
        print(f"  Audio: {Path(aud_path).name}")
        
        try:
            # Load audio to get duration
            audio = AudioFileClip(aud_path)
            duration = audio.duration
            print(f"Audio duration: {duration:.2f} seconds")
            
            # Create image clip
            image = ImageClip(img_path, duration=duration)
            
            # Set FPS using with_fps instead of set_fps
            image = image.with_fps(24)
            
            # Add audio using with_audio instead of set_audio  
            video_clip = image.with_audio(audio)
            
            clips.append(video_clip)
            
            print(f"Created clip successfully")
            
        except Exception as e:
            print(f"Error creating clip: {e}")
            continue
    
    if not clips:
        print("No clips created successfully!")
        st.session_state.progress_logs.error("No clips created successfully!")
        return
    
    print(f"\nConcatenating {len(clips)} clips...")
    
    try:
        # Concatenate all clips
        final_video = concatenate_videoclips(clips)
        
        print(f"Writing video to {st.session_state.video_file_path}...")
        
        # Write video file
        final_video.write_videofile(
            st.session_state.video_file_path,
            fps=24,
            codec='libx264',
            audio_codec='aac',
            logger=None
        )
        
        print("Video created successfully!")
        st.session_state.progress_logs.success("Video created successfully!")
        st.markdown("### Preview")
        st.video(st.session_state.video_file_path)

        try:
            with open(st.session_state.video_file_path, "rb") as file:
                video_bytes = file.read()

            st.download_button(
                label="Download Video",
                data=video_bytes,
                file_name="rcb_generated_video.mp4",
                mime="video/mp4"
            )
        except FileNotFoundError:
            st.error(f"Error: The video file '{video_file_path}' was not found.")
        except Exception as e:
            st.error(f"An error occurred: {e}")

        st.markdown("Note:: When you click the Download button, the video will disappear from the page. This is normal behavior. You can still save the file to your device.")


    except Exception as e:
        print(f"✗ Error creating final video: {e}")
    
    finally:
        # Clean up
        try:
            final_video.close()
            for clip in clips:
                clip.close()
        except:
            pass

def generate_video_from_pptx(uploaded_file):
    # Process the uploaded PPTX file and generate a video
    with open("input.pptx", "wb") as f:
        f.write(uploaded_file.getbuffer())

    # Convert pptx to pdf
    st.session_state.progress_logs.info("Converting PPTX to PDF...")
    subprocess.run([
        "libreoffice", "--headless", "--convert-to", "pdf",
        "input.pptx"
    ])

    # Convert PDF pages to images
    slides = convert_from_path("input.pdf", dpi=300)
    for i, slide in enumerate(slides, start=1):
        st.session_state.progress_logs.info(f"Saving slide {i} as image...")
        slide.save(f"{st.session_state.video_data_dir}/{i}.png", "PNG")

    ## Extract speaker notes
    def extract_notes(pptx_file):
        prs = Presentation(pptx_file)
        notes = []
        for i, slide in enumerate(prs.slides, start=1):
            text = ""
            if slide.has_notes_slide:
                text_frame = slide.notes_slide.notes_text_frame
                if text_frame:
                    text = "\n".join([p.text for p in text_frame.paragraphs if p.text])
            notes.append((i, text))
        return notes

    for num, note in extract_notes("input.pptx"):
        st.session_state.progress_logs.info(f"Extracting notes from slide {num}...")
        with open(f"{st.session_state.video_data_dir}/{num}.txt", "w") as f:
            f.write(note)
        st.session_state.curated_transcript = note
        st.session_state.audio_file_path_txt = f"{st.session_state.video_data_dir}/{num}.txt"
        st.session_state.audio_file_path_wav = f"{st.session_state.video_data_dir}/{num}.wav"
        st.session_state.audio_file_path_mp3 = f"{st.session_state.video_data_dir}/{num}.mp3"
        st.session_state.progress_logs.info(f"Creating audio file for slide {num} from transcript")
        print(f"Creating audio file for slide {num} from transcript")
        generate_audio_file_from_transcript()
        shutil.copyfile(st.session_state.default_audio_file_path_wav, f"{st.session_state.video_data_dir}/{num}.wav")
        shutil.copyfile(st.session_state.default_audio_file_path_mp3, f"{st.session_state.video_data_dir}/{num}.mp3")
        shutil.copyfile(st.session_state.default_audio_file_path_txt, f"{st.session_state.video_data_dir}/{num}.txt")

    st.session_state.progress_logs.info("Slides and notes extracted successfully!")
    simple_video_creator()

    #def combine_audio_files_and_images():
url = "https://docs.google.com/presentation/u/0/?ftv=1&tgif=d"

instructions = st.empty()
with instructions.container():
    st.markdown("### Template slide deck for video generation")
    st.markdown("Use any presentation template from [here](%s) for video generation " % url)
    st.markdown("Create copy of selected template and add your content with speaker notes.")
    st.markdown("Download the slide in .pptx format and use it for generating video")

# Sidebar on streamlit app
# with st.sidebar:
st.sidebar.subheader("Upload your presentation (.pptx) with speaker notes")
uploaded_file = st.sidebar.file_uploader(
    "Upload Presentation,",
    type=['pptx'],
    accept_multiple_files=False,
    help="Upload your presentation file (make sure it includes speaker notes).",
    disabled=st.session_state.disable_all
)

if uploaded_file:
    process_file_button = st.sidebar.button("Generate Video")
    if process_file_button:
        print(f"UPLOADED FILES:\n {uploaded_file}")
        instructions.empty()
        clear_old_files()
        generate_video_from_pptx(uploaded_file)

